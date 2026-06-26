import os
import io
import sys
from flask import Flask, render_template, request, jsonify, Response, stream_with_context
from google import genai
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image

sys.stdout.reconfigure(encoding='utf-8')

# Load environment variables
load_dotenv()

app = Flask(__name__, static_folder='static', template_folder='templates')
app.secret_key = os.getenv('SECRET_KEY', 'dev-key-change-in-production')

# Configure Google Gemini Client
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
if not GEMINI_API_KEY or "your_api_key" in GEMINI_API_KEY:
    print("[ERROR] GEMINI_API_KEY is missing or invalid in your .env file.")
    client = None
else:
    client = genai.Client(api_key=GEMINI_API_KEY)

# Helper to find a working model
def get_working_model():
    if not client: return "gemini-1.5-flash"
    try:
        print("[DEBUG] Attempting to list models...", flush=True)
        models = list(client.models.list())
        print(f"[DEBUG] Found {len(models)} models available.", flush=True)
        
        # Priority 1: Exact match for 1.5-flash
        for m in models:
            if m.name == "models/gemini-1.5-flash":
                return m.name
        
        # Priority 2: Any 1.5-flash variation
        for m in models:
            if m.name and "gemini-1.5-flash" in m.name:
                return m.name

        # Priority 3: First available model
        if models:
            return models[0].name
            
    except Exception as e:
        print(f"DEBUG: Error listing models: {e}")
    
    return "gemini-1.5-flash"

# Initialize model name
MODEL_NAME = get_working_model()
print(f"[INFO] Detected Working Model: {MODEL_NAME}")

def extract_text(file):
    filename = file.filename.lower()
    try:
        if filename.endswith(".pdf"):
            reader = PdfReader(file)
            return "".join([page.extract_text() for page in reader.pages])
        elif filename.endswith(".docx"):
            doc = Document(file)
            return "\n".join([p.text for p in doc.paragraphs])
        elif filename.endswith(".pptx"):
            prs = Presentation(file)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
    except Exception as e:
        print(f"Extraction error: {e}")
    return None

def process_multimodal(prompt, file):
    if not client: return "API Key is invalid or missing in .env"
    
    try:
        current_model = get_working_model()
        if file.filename.lower().endswith(('.png', '.jpg', '.jpeg', '.webp', '.heic', '.heif')):
            img = Image.open(file)
            response = client.models.generate_content(model=current_model, contents=[prompt, img])
            return response.text
        else:
            text_content = extract_text(file)
            if text_content:
                response = client.models.generate_content(model=current_model, contents=f"{prompt}\n\nContent:\n{text_content}")
                return response.text
            return "Could not extract text from file."
    except Exception as e:
        return f"Gemini Error: {str(e)}"

@app.route("/models")
def list_models():
    if not client: return jsonify({"error": "API Key not configured"}), 500
    try:
        models = []
        for m in client.models.list():
            models.append({"name": m.name})
        return jsonify({
            "detected_model": get_working_model(),
            "available_models": models
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/")
def index():
    return render_template("index.html")

@app.route("/ask", methods=["POST"])
def ask():
    if not client: return jsonify({"error": "API Key not configured"}), 500
    data = request.json
    question = data.get("question")
    lang = data.get("lang", "English")
    if not question:
        return jsonify({"error": "Question is required"}), 400
    try:
        current_model = get_working_model()
        prompt = f"You are EduGenie, an expert educational learning assistant. Answer this student question clearly in {lang}: {question}"
        response = client.models.generate_content(model=current_model, contents=prompt)
        return jsonify({"answer": response.text})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/chat", methods=['POST'])
def chat():
    """Chat endpoint with streaming response."""
    if not client:
        return jsonify({
            'error': 'API Key not configured',
            'status': 'error'
        }), 500
    
    data = request.get_json(silent=True) or {}
    message = (data.get('message') or '').strip()
    
    if not message:
        return jsonify({'error': 'empty message', 'status': 'error'}), 400
    
    if len(message) > 5000:
        return jsonify({'error': 'message too long (max 5000 chars)', 'status': 'error'}), 400

    def generate_stream():
        try:
            current_model = get_working_model()
            prompt = (
                "You are EduGenie, an expert educational learning assistant. "
                "Answer this student question clearly and comprehensively: " + message
            )
            # Stream the response from Gemini
            response = client.models.generate_content(
                model=current_model, 
                contents=prompt,
                #stream=True
            )
            
            # Yield each chunk as it arrives
            for chunk in response:
                if chunk.text:
                    yield f"data: {chunk.text}\n\n"
                    
        except Exception as e:
            yield f"data: [ERROR] {str(e)[:200]}\n\n"
    
    return Response(
        stream_with_context(generate_stream()),
        mimetype='text/event-stream',
        headers={
            'Cache-Control': 'no-cache',
            'X-Accel-Buffering': 'no',
            'Connection': 'keep-alive'
        }
    )

@app.route('/clear', methods=['POST'])
def clear_history():
    """Clear conversation history."""
    return jsonify({'status': 'cleared'})

@app.route("/quiz", methods=["POST"])
def quiz():
    if not client: return jsonify({"error": "API Key not configured"}), 500
    topic = request.form.get("topic")
    file = request.files.get("file")
    lang = request.form.get("lang", "English")
    
    schema = f"""
    Format your response as a JSON array of objects in {lang}. Each object must have:
    - "question": text of the question in {lang}
    - "options": {{"A": "...", "B": "...", "C": "...", "D": "..."}} (values in {lang})
    - "correct": the letter of the correct answer (A, B, C, or D)
    - "explanation": a brief explanation in {lang} of why that answer is correct.
    Return ONLY the JSON. No other text.
    """
    
    prompt = f"Generate a 5-question multiple-choice quiz. {schema}"
    
    try:
        current_model = get_working_model()
        if file:
            result_text = process_multimodal(f"{prompt} based on this content:", file)
        elif topic:
            response = client.models.generate_content(model=current_model, contents=f"{prompt} about {topic}")
            result_text = response.text
        else:
            return jsonify({"error": "Topic or File is required"}), 400

        import json
        cleaned_json = result_text.strip()
        if cleaned_json.startswith("```json"):
            cleaned_json = cleaned_json[7:-3].strip()
        elif cleaned_json.startswith("```"):
            cleaned_json = cleaned_json[3:-3].strip()
        
        quiz_data = json.loads(cleaned_json)
        return jsonify({"quiz": quiz_data})
        
    except Exception as e:
        print(f"Quiz generation error: {e}")
        return jsonify({"error": "Failed to generate structured quiz in the selected language. Try again."}), 500

@app.route("/summary", methods=["POST"])
def summary():
    if not client: return jsonify({"error": "API Key not configured"}), 500
    content = request.form.get("content")
    file = request.files.get("file")
    lang = request.form.get("lang", "English")
    
    prompt = f"Provide a concise and structured summary in {lang} of the following material:"
    
    try:
        current_model = get_working_model()
        if file:
            result = process_multimodal(prompt, file)
            return jsonify({"summary": result})
        elif content:
            response = client.models.generate_content(model=current_model, contents=f"{prompt}\n\n{content}")
            return jsonify({"summary": response.text})
        return jsonify({"error": "Content or File is required"}), 400
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/health', methods=['GET'])
def health():
    """Simple health check."""
    return jsonify({'status': 'ok', 'service': 'EduGenie'})

if __name__ == "__main__":
    print("[INFO] EduGenie is running on http://127.0.0.1:5000")
    app.run(debug=True)
